"""Parse a .pptx file into a Presentation using python-pptx."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation as PptxPresentation

from pitchdeck.models import Presentation, Slide


def load_pptx(path: Path) -> Presentation:
	"""Open a .pptx file and extract slide data into a Presentation.

	Raises FileNotFoundError if path does not exist.
	Raises ValueError if the file cannot be parsed as a valid pptx.
	"""
	if not path.exists():
		raise FileNotFoundError(f"Presentation file not found: {path}")
	try:
		prs = PptxPresentation(str(path))
	except Exception as exc:
		raise ValueError(f"Could not parse {path.name} as a PowerPoint file: {exc}") from exc

	slides = [_extract_slide(slide, i + 1) for i, slide in enumerate(prs.slides)]
	return Presentation(slides=slides, source_path=str(path.resolve()), source_format="pptx")


def _extract_slide(slide: Any, number: int) -> Slide:
	"""Extract structured data from a single python-pptx Slide object."""
	return Slide(
		number=number,
		title=_get_title(slide),
		body=_get_body_texts(slide),
		notes=_get_notes(slide),
	)


def _get_title(slide: Any) -> str:
	"""Return the title placeholder text, or empty string if absent."""
	title_shape = slide.shapes.title
	if title_shape is None:
		return ""
	text = title_shape.text_frame.text if title_shape.has_text_frame else ""
	return text.strip()


def _get_body_texts(slide: Any) -> list[str]:
	"""Return non-empty text strings from all non-title text frames."""
	title_shape = slide.shapes.title
	texts: list[str] = []
	for shape in slide.shapes:
		if shape is title_shape:
			continue
		if not shape.has_text_frame:
			continue
		for para in shape.text_frame.paragraphs:
			text = para.text.strip()
			if text:
				texts.append(text)
	return texts


def _get_notes(slide: Any) -> str:
	"""Return speaker notes text, or empty string if none exist."""
	try:
		# notes_slide raises AttributeError on some pptx versions when no notes exist
		notes_tf = slide.notes_slide.notes_text_frame
		return notes_tf.text.strip()
	except (AttributeError, TypeError, ValueError):
		return ""
